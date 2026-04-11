"""StudyCoach API routes: courses, materials, sections, study plan, teaching mode."""
import uuid
import re
from datetime import datetime, timezone
from typing import Optional

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, desc
from pydantic import BaseModel

from app.core.database import get_db
from app.core.auth import hash_password, verify_password, create_access_token
from app.core.config import get_settings
from app.models.models import (
    StudyCoachUser, SCCourse, SCMaterial, SCSection,
    SCPlanItem, SCTeachSession, SCTeachMessage,
)
from app.services.ai_providers import get_llm_provider

settings = get_settings()
router = APIRouter(prefix="/study")

SBU_DOMAINS = {"stonybrook.edu", "alumni.stonybrook.edu", "cs.stonybrook.edu"}

# ─── Auth dependency ──────────────────────────────────────────────────

def _get_user_dep():
    from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
    from fastapi import Security
    from jose import jwt as jose_jwt

    bearer = HTTPBearer(auto_error=False)

    async def dep(
        credentials: HTTPAuthorizationCredentials = Security(bearer),
        db: AsyncSession = Depends(get_db),
    ) -> StudyCoachUser:
        if not credentials:
            raise HTTPException(status_code=401, detail="Not authenticated")
        try:
            payload = jose_jwt.decode(
                credentials.credentials, settings.jwt_secret,
                algorithms=[settings.jwt_algorithm],
            )
            user_id = payload.get("sub")
            if not user_id:
                raise HTTPException(status_code=401, detail="Invalid token")
        except Exception:
            raise HTTPException(status_code=401, detail="Invalid token")
        result = await db.execute(select(StudyCoachUser).where(StudyCoachUser.id == uuid.UUID(user_id)))
        user = result.scalar_one_or_none()
        if not user or not user.is_active:
            raise HTTPException(status_code=401, detail="User not found")
        return user

    return dep

get_current_user = _get_user_dep()


# ─── Schemas ──────────────────────────────────────────────────────────

class RegisterRequest(BaseModel):
    email: str
    password: str
    name: str

class LoginRequest(BaseModel):
    email: str
    password: str

class CourseCreate(BaseModel):
    code: str
    name: str
    description: Optional[str] = None

class PlanItemCreate(BaseModel):
    title: str
    due_date: Optional[datetime] = None
    item_type: str = "study"
    notes: Optional[str] = None

class PlanItemUpdate(BaseModel):
    title: Optional[str] = None
    due_date: Optional[datetime] = None
    item_type: Optional[str] = None
    is_completed: Optional[bool] = None
    notes: Optional[str] = None

class TeachRequest(BaseModel):
    course_id: str
    message: str
    session_id: Optional[str] = None
    section_id: Optional[str] = None
    knowledge_level: Optional[str] = None  # unknown/partial/confident


# ─── File parsing ─────────────────────────────────────────────────────

CODE_EXTENSIONS = {".py", ".js", ".ts", ".java", ".cpp", ".c", ".cs", ".go", ".rb", ".r", ".m", ".swift"}

def _detect_type(filename: str) -> str:
    fn = filename.lower()
    if fn.endswith(".pdf"):        return "pdf"
    if fn.endswith(".docx"):       return "docx"
    if fn.endswith(".pptx"):       return "pptx"
    if fn.endswith(".csv"):        return "csv"
    if fn.endswith((".txt", ".md")): return "txt"
    ext = "." + fn.rsplit(".", 1)[-1] if "." in fn else ""
    if ext in CODE_EXTENSIONS:    return "code"
    return "txt"

def _sanitize(text: str) -> str:
    """Remove null bytes and other chars PostgreSQL rejects in UTF-8."""
    return text.replace("\x00", "").strip()


def _parse_file(content: bytes, file_type: str, filename: str) -> list[dict]:
    """Parse file into raw page/slide chunks: [{"page": int, "content": str}]"""
    try:
        if file_type == "pdf":
            import pdfplumber, io
            pages = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for i, page in enumerate(pdf.pages):
                    t = page.extract_text()
                    if t and t.strip():
                        pages.append({"page": i + 1, "content": _sanitize(t)})
            return pages

        if file_type == "docx":
            import docx, io
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return [{"page": 1, "content": _sanitize(text)}]

        if file_type == "pptx":
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            pages = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    pages.append({"page": i, "content": _sanitize("\n".join(texts))})
            return pages

        if file_type == "csv":
            import pandas as pd, io
            df = pd.read_csv(io.BytesIO(content))
            summary = f"Columns: {', '.join(df.columns)}\nRows: {len(df)}\n\nFirst 10 rows:\n{df.head(10).to_string()}"
            return [{"page": 1, "content": _sanitize(summary)}]

        # txt, code, fallback
        text = content.decode("utf-8", errors="ignore")
        return [{"page": 1, "content": _sanitize(text)}]

    except Exception:
        text = content.decode("utf-8", errors="ignore")
        return [{"page": 1, "content": text}]


def _group_pages(pages: list[dict], file_type: str) -> list[dict]:
    """
    Group pages into logical sections.
    - pptx: each slide is a section
    - pdf/docx/txt: group into ~600-word chunks, respecting page boundaries
    """
    if not pages:
        return []

    if file_type == "pptx":
        return [{"page_start": p["page"], "content": p["content"]} for p in pages]

    # For pdf/docx: merge pages into ~600-word chunks
    sections = []
    current_content = ""
    current_start = pages[0]["page"]
    target_words = 600

    for p in pages:
        text = p["content"]
        # If adding this page stays under limit, append; otherwise flush
        combined = (current_content + "\n\n" + text).strip()
        if current_content and len(combined.split()) > target_words:
            sections.append({"page_start": current_start, "content": current_content.strip()})
            current_content = text
            current_start = p["page"]
        else:
            current_content = combined

    if current_content.strip():
        sections.append({"page_start": current_start, "content": current_content.strip()})

    return sections


async def _ai_enrich_sections(raw_groups: list[dict], course_name: str, filename: str) -> list[dict]:
    """Use LLM to generate title + difficulty + concepts + prerequisites for each section."""
    import json as _json
    llm = get_llm_provider()
    enriched = []

    for i, group in enumerate(raw_groups):
        content_preview = group["content"][:1200]
        prompt = f"""You are analyzing a section of a university course document.

Course: "{course_name}"
File: "{filename}"
Content:
{content_preview}

Return ONLY valid JSON (no markdown, no extra text):
{{
  "title": "<concise descriptive title for this section, max 8 words>",
  "difficulty": <integer 1-5, where 1=intro/easy 5=very advanced>,
  "concepts": ["key concept 1", "key concept 2", "key concept 3"],
  "prerequisites": ["prereq 1", "prereq 2"]
}}

For title: be specific and descriptive (e.g., "Measures of Central Tendency", "Linear Regression Basics", "R Data Structures").
Do NOT use page numbers or author names as titles."""

        try:
            raw = await llm.generate(prompt, system="You are a curriculum analyst. Return only valid JSON.")
            match = re.search(r"\{.*\}", raw, re.DOTALL)
            data = _json.loads(match.group()) if match else {}

            title = str(data.get("title") or f"Section {i + 1}").strip()
            # Sanitize: reject titles that look like raw content (too short, numbers, etc.)
            if len(title) < 4 or title[0].isdigit() and len(title) < 8:
                title = f"Section {i + 1}"

            enriched.append({
                "title": title,
                "content": group["content"],
                "page_start": group.get("page_start", i + 1),
                "difficulty": max(1, min(5, int(data.get("difficulty") or 3))),
                "concepts": (data.get("concepts") or [])[:6],
                "prerequisites": (data.get("prerequisites") or [])[:4],
                "order": i,
            })
        except Exception:
            enriched.append({
                "title": f"Section {i + 1}",
                "content": group["content"],
                "page_start": group.get("page_start", i + 1),
                "difficulty": 3,
                "concepts": [],
                "prerequisites": [],
                "order": i,
            })

    return enriched


# ─── Teaching system ──────────────────────────────────────────────────

TEACH_SYSTEM = """You are StudyCoach, a Socratic tutor for Stony Brook University students.

CONVERSATION STYLE — follow this strictly:
- Give ONE idea, concept, or hint per message. No more.
- Keep each response to 3–5 sentences unless a code snippet is truly necessary.
- End every response with exactly ONE question to check understanding or prompt the student to try.
- Do NOT give multi-part explanations or numbered breakdowns in a single message.
- Do NOT dump everything you know upfront. Teach like a conversation, not a lecture.

FORMATTING RULES:
- Never use headers (no #, ##, ###).
- Use **bold** only for key terms, sparingly (1–2 per message max).
- Use code blocks for code or technical syntax only.
- Plain conversational sentences otherwise — no bullet lists of explanations.

ABSOLUTE RULES:
1. NEVER give the direct answer to a problem or assignment.
2. If asked "just tell me the answer," say: "Let's figure it out together — what do you think happens when...?"
3. If the student is wrong, don't just correct them — ask what made them think that first.
4. After the student understands one thing, naturally move to the next piece.
5. Adapt depth to knowledge level: unknown = analogy first; partial = skip basics; confident = challenge with edge cases.
6. Ground explanations in the provided course material when available."""


def _build_teach_prompt(
    message: str,
    history: list[dict],
    knowledge_level: str,
    section_content: str | None,
    course_name: str,
) -> str:
    context = ""
    if section_content:
        context = f"\n[COURSE MATERIAL]\n{section_content[:2000]}\n[/COURSE MATERIAL]\n"

    history_block = ""
    if history:
        lines = [f"{'Student' if m['role'] == 'user' else 'Coach'}: {m['content']}" for m in history[-8:]]
        history_block = "\n[CONVERSATION]\n" + "\n".join(lines) + "\n[/CONVERSATION]\n"

    level_note = {
        "unknown": "The student is new to this topic. Start from the very basics.",
        "partial": "The student has heard of this but is fuzzy on the details.",
        "confident": "The student knows the basics — skip fundamentals and focus on application.",
    }.get(knowledge_level, "")

    return f"""Course: {course_name}
Student knowledge level: {knowledge_level} — {level_note}
{context}{history_block}
Student: {message}"""


# ─── Auth routes ──────────────────────────────────────────────────────

@router.post("/auth/register", status_code=201)
async def register(req: RegisterRequest, db: AsyncSession = Depends(get_db)):
    domain = req.email.lower().split("@")[-1]
    if domain not in SBU_DOMAINS:
        raise HTTPException(status_code=400, detail="Must be an @stonybrook.edu email")
    existing = (await db.execute(select(StudyCoachUser).where(StudyCoachUser.email == req.email))).scalar_one_or_none()
    if existing:
        raise HTTPException(status_code=409, detail="Email already registered")
    user = StudyCoachUser(email=req.email, password_hash=hash_password(req.password), name=req.name)
    db.add(user)
    await db.flush()
    token = create_access_token({"sub": str(user.id), "email": user.email, "name": user.name, "type": "study"})
    return {"access_token": token, "user_id": str(user.id), "name": user.name, "email": user.email}


@router.post("/auth/login")
async def login(req: LoginRequest, db: AsyncSession = Depends(get_db)):
    domain = req.email.lower().split("@")[-1]
    if domain not in SBU_DOMAINS:
        raise HTTPException(status_code=400, detail="Must be an @stonybrook.edu email")
    result = await db.execute(select(StudyCoachUser).where(StudyCoachUser.email == req.email, StudyCoachUser.is_active == True))
    user = result.scalar_one_or_none()
    if not user or not verify_password(req.password, user.password_hash):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    user.last_login_at = datetime.now(timezone.utc)
    token = create_access_token({"sub": str(user.id), "email": user.email, "name": user.name, "type": "study"})
    return {"access_token": token, "user_id": str(user.id), "name": user.name, "email": user.email}


# ─── Course routes ─────────────────────────────────────────────────────

@router.post("/courses", status_code=201)
async def create_course(
    req: CourseCreate,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    course = SCCourse(user_id=user.id, code=req.code.upper(), name=req.name, description=req.description)
    db.add(course)
    await db.flush()
    return _course_dict(course)


@router.get("/courses")
async def list_courses(user: StudyCoachUser = Depends(get_current_user), db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(SCCourse).where(SCCourse.user_id == user.id).order_by(desc(SCCourse.created_at))
    )
    return [_course_dict(c) for c in result.scalars().all()]


@router.get("/courses/{course_id}")
async def get_course(
    course_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    course = await _get_course(course_id, user.id, db)
    materials = (await db.execute(
        select(SCMaterial).where(SCMaterial.course_id == course_id).order_by(SCMaterial.created_at)
    )).scalars().all()
    # sections grouped by material_id for the frontend tabs
    sections_by_material: dict[str, list] = {}
    all_sections = []
    for mat in materials:
        secs = (await db.execute(
            select(SCSection).where(SCSection.material_id == mat.id).order_by(SCSection.order)
        )).scalars().all()
        mat_secs = [_section_dict(s, mat.filename) for s in secs]
        sections_by_material[str(mat.id)] = mat_secs
        all_sections.extend(mat_secs)
    return {
        **_course_dict(course),
        "materials": [_material_dict(m) for m in materials],
        "sections": all_sections,
        "sections_by_material": sections_by_material,
    }


@router.delete("/courses/{course_id}", status_code=204)
async def delete_course(
    course_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    course = await _get_course(course_id, user.id, db)
    await db.delete(course)


# ─── Material / Upload routes ──────────────────────────────────────────

@router.post("/courses/{course_id}/upload")
async def upload_material(
    course_id: uuid.UUID,
    file: UploadFile = File(...),
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    course = await _get_course(course_id, user.id, db)
    content = await file.read()
    fname = file.filename or "upload"
    file_type = _detect_type(fname)

    material = SCMaterial(
        course_id=course_id,
        filename=fname,
        file_type=file_type,
        status="processing",
    )
    db.add(material)
    await db.flush()

    # Parse into page chunks
    pages = _parse_file(content, file_type, fname)
    raw_text = _sanitize("\n\n".join(p["content"] for p in pages))
    material.raw_text = raw_text

    # Group into logical sections
    groups = _group_pages(pages, file_type)[:20]  # cap at 20 sections

    # AI enrich: title + difficulty + concepts + prerequisites
    enriched = await _ai_enrich_sections(groups, course.name, fname)

    for sec_data in enriched:
        sec = SCSection(
            material_id=material.id,
            title=sec_data["title"],
            order=sec_data["order"],
            content=sec_data["content"],
            difficulty=sec_data["difficulty"],
            concepts=sec_data["concepts"],
            prerequisites=sec_data["prerequisites"],
        )
        db.add(sec)

    material.status = "ready"
    material.parsed_at = datetime.now(timezone.utc)

    await db.flush()
    return _material_dict(material)


@router.delete("/sections/{section_id}", status_code=204)
async def delete_section(
    section_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    sec = (await db.execute(select(SCSection).where(SCSection.id == section_id))).scalar_one_or_none()
    if not sec:
        raise HTTPException(status_code=404, detail="Section not found")
    # Verify ownership through material → course
    mat = (await db.execute(select(SCMaterial).where(SCMaterial.id == sec.material_id))).scalar_one_or_none()
    if mat:
        await _get_course(mat.course_id, user.id, db)
    await db.delete(sec)


@router.delete("/courses/{course_id}/materials/{material_id}", status_code=204)
async def delete_material(
    course_id: uuid.UUID,
    material_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    await _get_course(course_id, user.id, db)
    mat = (await db.execute(select(SCMaterial).where(SCMaterial.id == material_id, SCMaterial.course_id == course_id))).scalar_one_or_none()
    if not mat:
        raise HTTPException(status_code=404, detail="Material not found")
    await db.delete(mat)


# ─── Study Plan routes ─────────────────────────────────────────────────

@router.get("/courses/{course_id}/plan")
async def get_plan(
    course_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    await _get_course(course_id, user.id, db)
    result = await db.execute(
        select(SCPlanItem).where(SCPlanItem.course_id == course_id).order_by(SCPlanItem.due_date.nullslast(), SCPlanItem.created_at)
    )
    return [_plan_dict(p) for p in result.scalars().all()]


@router.post("/courses/{course_id}/plan", status_code=201)
async def add_plan_item(
    course_id: uuid.UUID,
    req: PlanItemCreate,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    await _get_course(course_id, user.id, db)
    item = SCPlanItem(course_id=course_id, title=req.title, due_date=req.due_date,
                      item_type=req.item_type, notes=req.notes)
    db.add(item)
    await db.flush()
    return _plan_dict(item)


@router.put("/plan/{item_id}")
async def update_plan_item(
    item_id: uuid.UUID,
    req: PlanItemUpdate,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    item = (await db.execute(select(SCPlanItem).where(SCPlanItem.id == item_id))).scalar_one_or_none()
    if not item:
        raise HTTPException(status_code=404, detail="Plan item not found")
    # Verify ownership
    course = (await db.execute(select(SCCourse).where(SCCourse.id == item.course_id, SCCourse.user_id == user.id))).scalar_one_or_none()
    if not course:
        raise HTTPException(status_code=403, detail="Not your plan item")
    for field, value in req.model_dump(exclude_unset=True).items():
        setattr(item, field, value)
    return _plan_dict(item)


@router.delete("/plan/{item_id}", status_code=204)
async def delete_plan_item(
    item_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    item = (await db.execute(select(SCPlanItem).where(SCPlanItem.id == item_id))).scalar_one_or_none()
    if not item:
        raise HTTPException(status_code=404, detail="Plan item not found")
    course = (await db.execute(select(SCCourse).where(SCCourse.id == item.course_id, SCCourse.user_id == user.id))).scalar_one_or_none()
    if not course:
        raise HTTPException(status_code=403, detail="Not your plan item")
    await db.delete(item)


# ─── AI Plan generation from syllabus ─────────────────────────────────

@router.post("/courses/{course_id}/plan/generate")
async def generate_plan(
    course_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Parse syllabus content from uploaded materials and generate initial plan items."""
    course = await _get_course(course_id, user.id, db)

    # Get all material text
    materials = (await db.execute(
        select(SCMaterial).where(SCMaterial.course_id == course_id, SCMaterial.status == "parsed")
    )).scalars().all()

    if not materials:
        raise HTTPException(status_code=400, detail="Upload course materials first")

    combined = "\n\n".join(
        f"[{m.filename}]\n{(m.raw_text or '')[:2000]}" for m in materials[:3]
    )

    llm = get_llm_provider()
    prompt = f"""Extract study plan items from these course materials for "{course.name}".

{combined}

Return ONLY a JSON array of plan items. Extract deadlines, exams, assignments, and weekly topics:
[
  {{"title": "Read Chapter 1", "due_date": null, "item_type": "lecture", "notes": ""}},
  {{"title": "Assignment 1 due", "due_date": "2026-04-20T23:59:00", "item_type": "assignment", "notes": ""}},
  {{"title": "Midterm Exam", "due_date": "2026-04-25T10:00:00", "item_type": "exam", "notes": ""}}
]

Use null for due_date if no date found. item_type: lecture/assignment/exam/review/study. Return 5-15 items."""

    try:
        raw = await llm.generate(prompt, system="Return only valid JSON array. No markdown.")
        import json
        match = re.search(r"\[.*\]", raw, re.DOTALL)
        if not match:
            raise ValueError("No JSON array found")
        items_data = json.loads(match.group())

        created = []
        for item_data in items_data[:15]:
            due = None
            if item_data.get("due_date"):
                try:
                    due = datetime.fromisoformat(item_data["due_date"].replace("Z", "+00:00"))
                except Exception:
                    due = None
            item = SCPlanItem(
                course_id=course_id,
                title=item_data.get("title", "Study session"),
                due_date=due,
                item_type=item_data.get("item_type", "study"),
                notes=item_data.get("notes", "") or None,
            )
            db.add(item)
            created.append(item)
        await db.flush()
        return [_plan_dict(p) for p in created]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Plan generation failed: {str(e)}")


# ─── Teaching mode ─────────────────────────────────────────────────────

@router.post("/teach")
async def teach(
    req: TeachRequest,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    course = await _get_course(uuid.UUID(req.course_id), user.id, db)

    # Get or create session
    session = None
    if req.session_id:
        session = (await db.execute(
            select(SCTeachSession).where(SCTeachSession.id == uuid.UUID(req.session_id), SCTeachSession.user_id == user.id)
        )).scalar_one_or_none()

    if not session:
        section_id = uuid.UUID(req.section_id) if req.section_id else None
        session = SCTeachSession(
            course_id=uuid.UUID(req.course_id),
            user_id=user.id,
            section_id=section_id,
            knowledge_level=req.knowledge_level or "unknown",
        )
        db.add(session)
        await db.flush()

    if req.knowledge_level and req.knowledge_level != session.knowledge_level:
        session.knowledge_level = req.knowledge_level

    # Fetch section content
    section_content = None
    if session.section_id:
        sec = (await db.execute(select(SCSection).where(SCSection.id == session.section_id))).scalar_one_or_none()
        if sec:
            section_content = sec.content

    # Fetch history
    history_rows = (await db.execute(
        select(SCTeachMessage).where(SCTeachMessage.session_id == session.id).order_by(SCTeachMessage.created_at.desc()).limit(10)
    )).scalars().all()
    history = [{"role": m.role, "content": m.content} for m in reversed(history_rows)]

    # Build prompt and call LLM
    prompt = _build_teach_prompt(req.message, history, session.knowledge_level, section_content, course.name)
    llm = get_llm_provider()
    response = await llm.generate(prompt, system=TEACH_SYSTEM)

    # Save messages
    db.add(SCTeachMessage(session_id=session.id, role="user", content=req.message))
    db.add(SCTeachMessage(session_id=session.id, role="assistant", content=response))
    await db.flush()

    return {"response": response, "session_id": str(session.id)}


@router.get("/teach/{session_id}/messages")
async def get_teach_messages(
    session_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    session = (await db.execute(
        select(SCTeachSession).where(SCTeachSession.id == session_id, SCTeachSession.user_id == user.id)
    )).scalar_one_or_none()
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    messages = (await db.execute(
        select(SCTeachMessage).where(SCTeachMessage.session_id == session_id).order_by(SCTeachMessage.created_at)
    )).scalars().all()
    return [{"id": str(m.id), "role": m.role, "content": m.content, "created_at": m.created_at.isoformat()} for m in messages]


@router.get("/sessions")
async def list_all_sessions(
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """All teach sessions across all courses for the current user."""
    sessions = (await db.execute(
        select(SCTeachSession, SCCourse.name.label("course_name"), SCCourse.code.label("course_code"))
        .join(SCCourse, SCTeachSession.course_id == SCCourse.id)
        .where(SCTeachSession.user_id == user.id)
        .order_by(desc(SCTeachSession.created_at))
    )).all()

    result = []
    for row in sessions:
        s = row[0]
        msgs = (await db.execute(
            select(SCTeachMessage).where(SCTeachMessage.session_id == s.id).order_by(SCTeachMessage.created_at)
        )).scalars().all()
        preview = msgs[0].content[:120] if msgs else ""
        section_title = None
        if s.section_id:
            sec = (await db.execute(select(SCSection).where(SCSection.id == s.section_id))).scalar_one_or_none()
            section_title = sec.title if sec else None
        result.append({
            "id": str(s.id),
            "course_id": str(s.course_id),
            "course_name": row[1],
            "course_code": row[2],
            "knowledge_level": s.knowledge_level,
            "section_id": str(s.section_id) if s.section_id else None,
            "section_title": section_title,
            "message_count": len(msgs),
            "preview": preview,
            "created_at": s.created_at.isoformat(),
        })
    return result


@router.get("/courses/{course_id}/sessions")
async def list_sessions(
    course_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    await _get_course(course_id, user.id, db)
    sessions = (await db.execute(
        select(SCTeachSession)
        .where(SCTeachSession.course_id == course_id, SCTeachSession.user_id == user.id)
        .order_by(desc(SCTeachSession.created_at))
    )).scalars().all()

    result = []
    for s in sessions:
        msgs = (await db.execute(
            select(SCTeachMessage).where(SCTeachMessage.session_id == s.id).order_by(SCTeachMessage.created_at)
        )).scalars().all()
        preview = msgs[0].content[:120] if msgs else ""
        section_title = None
        if s.section_id:
            sec = (await db.execute(select(SCSection).where(SCSection.id == s.section_id))).scalar_one_or_none()
            section_title = sec.title if sec else None
        result.append({
            "id": str(s.id),
            "knowledge_level": s.knowledge_level,
            "section_id": str(s.section_id) if s.section_id else None,
            "section_title": section_title,
            "message_count": len(msgs),
            "preview": preview,
            "created_at": s.created_at.isoformat(),
        })
    return result


@router.delete("/sessions/{session_id}", status_code=204)
async def delete_session(
    session_id: uuid.UUID,
    user: StudyCoachUser = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    session = (await db.execute(
        select(SCTeachSession).where(SCTeachSession.id == session_id, SCTeachSession.user_id == user.id)
    )).scalar_one_or_none()
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    await db.delete(session)


# ─── Helpers ──────────────────────────────────────────────────────────

async def _get_course(course_id: uuid.UUID, user_id: uuid.UUID, db: AsyncSession) -> SCCourse:
    course = (await db.execute(
        select(SCCourse).where(SCCourse.id == course_id, SCCourse.user_id == user_id)
    )).scalar_one_or_none()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    return course


def _course_dict(c: SCCourse) -> dict:
    return {"id": str(c.id), "code": c.code, "name": c.name, "description": c.description, "created_at": c.created_at.isoformat()}


def _material_dict(m: SCMaterial) -> dict:
    return {"id": str(m.id), "filename": m.filename, "file_type": m.file_type, "status": m.status,
            "parsed_at": m.parsed_at.isoformat() if m.parsed_at else None, "created_at": m.created_at.isoformat()}


def _section_dict(s: SCSection, material_filename: str = "") -> dict:
    return {"id": str(s.id), "material_id": str(s.material_id), "material_filename": material_filename,
            "title": s.title, "order": s.order, "difficulty": s.difficulty,
            "concepts": s.concepts or [], "prerequisites": s.prerequisites or []}


def _plan_dict(p: SCPlanItem) -> dict:
    return {"id": str(p.id), "course_id": str(p.course_id), "title": p.title,
            "due_date": p.due_date.isoformat() if p.due_date else None,
            "item_type": p.item_type, "is_completed": p.is_completed, "notes": p.notes,
            "created_at": p.created_at.isoformat()}
